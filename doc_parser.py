#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文档解析与加载模块
基于 Karpathy 极简本地 RAG 知识库架构

功能：
1. 支持 .docx / .pdf / .pptx / .txt 格式文档解析
2. 自动提取文档正文、标题、列表、表格文本
3. 文本清洗与标准化处理
4. 批量文件夹遍历处理

使用方法：
    from doc_parser import extract_full_text, batch_load_docs
    
    # 单文件解析
    text = extract_full_text("document.pdf")
    
    # 批量文件夹解析
    docs = batch_load_docs("raw/")
"""

import os
import sys
import re
import logging
from pathlib import Path
from typing import Dict, List, Optional
from datetime import datetime

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# 支持的文件格式
SUPPORTED_FORMATS = {'.txt', '.docx', '.pdf', '.pptx'}


def clean_text(text: str) -> str:
    """
    文本清洗与标准化
    
    功能：
    1. 去除多余空行和空白字符
    2. 统一 UTF-8 编码
    3. 合并重复空白与换行
    4. 去除不可见特殊字符
    
    参数：
        text: 原始文本
    
    返回：
        清洗后的文本
    """
    if not text:
        return ""
    
    # 统一换行符
    text = text.replace('\r\n', '\n').replace('\r', '\n')
    
    # 去除不可见特殊字符（保留中文、英文、数字、常用标点）
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
    
    # 合并连续空行（最多保留2个换行）
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # 合并连续空格（保留单个空格）
    text = re.sub(r'[^\S\n]{2,}', ' ', text)
    
    # 去除每行首尾空白
    lines = [line.strip() for line in text.split('\n')]
    
    # 去除空行（保留段落间的单个空行）
    cleaned_lines = []
    prev_was_empty = False
    
    for line in lines:
        if not line:
            if not prev_was_empty:
                cleaned_lines.append('')
            prev_was_empty = True
        else:
            cleaned_lines.append(line)
            prev_was_empty = False
    
    # 去除开头和结尾的空行
    while cleaned_lines and not cleaned_lines[0]:
        cleaned_lines.pop(0)
    while cleaned_lines and not cleaned_lines[-1]:
        cleaned_lines.pop()
    
    return '\n'.join(cleaned_lines)


def detect_file_type(file_path: str) -> str:
    """
    检测文件类型
    
    参数：
        file_path: 文件路径
    
    返回：
        文件类型后缀（小写）
    """
    return Path(file_path).suffix.lower()


def is_supported_format(file_path: str) -> bool:
    """
    检查文件是否为支持的格式
    
    参数：
        file_path: 文件路径
    
    返回：
        是否支持
    """
    return detect_file_type(file_path) in SUPPORTED_FORMATS


def extract_txt(file_path: str) -> str:
    """
    提取纯文本文件内容
    
    参数：
        file_path: 文件路径
    
    返回：
        文本内容
    """
    # 尝试多种编码读取
    encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
    
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    
    # 如果所有编码都失败，使用二进制读取并忽略错误
    with open(file_path, 'rb') as f:
        return f.read().decode('utf-8', errors='ignore')


def extract_docx(file_path: str) -> str:
    """
    提取 Word 文档内容
    
    参数：
        file_path: 文件路径
    
    返回：
        提取的文本内容
    """
    try:
        from docx import Document
        
        doc = Document(file_path)
        text_parts = []
        
        # 提取段落文本
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        # 提取表格文本
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(' | '.join(row_text))
        
        return '\n\n'.join(text_parts)
        
    except Exception as e:
        logger.error(f"DOCX 解析失败: {e}")
        return ""


def extract_pdf(file_path: str) -> str:
    """
    提取 PDF 文档内容
    
    参数：
        file_path: 文件路径
    
    返回：
        提取的文本内容
    """
    try:
        from PyPDF2 import PdfReader
        
        reader = PdfReader(file_path)
        text_parts = []
        
        for page_num, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text_parts.append(f"[第 {page_num + 1} 页]\n{page_text}")
        
        return '\n\n'.join(text_parts)
        
    except Exception as e:
        logger.error(f"PDF 解析失败: {e}")
        return ""


def extract_pptx(file_path: str) -> str:
    """
    提取 PowerPoint 文档内容
    
    参数：
        file_path: 文件路径
    
    返回：
        提取的文本内容
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
                
                # 提取表格内容
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_texts.append(' | '.join(row_text))
            
            if slide_texts:
                text_parts.append(f"[第 {slide_num} 页]\n" + '\n'.join(slide_texts))
        
        return '\n\n'.join(text_parts)
        
    except Exception as e:
        logger.error(f"PPTX 解析失败: {e}")
        return ""


def extract_full_text(file_path: str) -> str:
    """
    提取单个文件的完整文本
    
    这是主要的公开接口，支持所有格式
    
    参数：
        file_path: 文件绝对路径或相对路径
    
    返回：
        清洗后的完整文本字符串
    
    使用示例：
        text = extract_full_text("raw/document.pdf")
        print(text)
    """
    file_path = str(Path(file_path).resolve())
    
    if not os.path.exists(file_path):
        logger.error(f"文件不存在: {file_path}")
        return ""
    
    file_name = os.path.basename(file_path)
    file_type = detect_file_type(file_path)
    
    logger.info(f"解析文件: {file_name} ({file_type})")
    
    try:
        # 根据文件类型选择解析方法
        if file_type == '.txt':
            raw_text = extract_txt(file_path)
        elif file_type == '.docx':
            raw_text = extract_docx(file_path)
        elif file_type == '.pdf':
            raw_text = extract_pdf(file_path)
        elif file_type == '.pptx':
            raw_text = extract_pptx(file_path)
        else:
            logger.warning(f"不支持的文件格式: {file_type}")
            return ""
        
        # 文本清洗
        cleaned_text = clean_text(raw_text)
        
        logger.info(f"解析完成: {file_name} ({len(cleaned_text)} 字符)")
        return cleaned_text
        
    except Exception as e:
        logger.error(f"解析文件失败 {file_name}: {e}")
        return ""


def batch_load_docs(folder_path: str) -> Dict[str, str]:
    """
    批量加载文件夹中的所有文档
    
    参数：
        folder_path: 文件夹路径
    
    返回：
        {文件绝对路径: 文本内容} 字典
    
    使用示例：
        docs = batch_load_docs("raw/")
        for path, content in docs.items():
            print(f"{path}: {len(content)} 字符")
    """
    folder_path = Path(folder_path).resolve()
    
    if not folder_path.exists():
        logger.error(f"文件夹不存在: {folder_path}")
        return {}
    
    if not folder_path.is_dir():
        logger.error(f"路径不是文件夹: {folder_path}")
        return {}
    
    results = {}
    total_files = 0
    success_files = 0
    failed_files = 0
    
    logger.info(f"开始扫描文件夹: {folder_path}")
    
    # 遍历所有文件和子文件夹
    for root, dirs, files in os.walk(folder_path):
        for file_name in files:
            file_path = Path(root) / file_name
            
            # 检查是否为支持的格式
            if not is_supported_format(str(file_path)):
                continue
            
            total_files += 1
            
            try:
                content = extract_full_text(str(file_path))
                
                if content:
                    results[str(file_path)] = content
                    success_files += 1
                    logger.info(f"✓ 成功: {file_name}")
                else:
                    failed_files += 1
                    logger.warning(f"✗ 内容为空: {file_name}")
                    
            except Exception as e:
                failed_files += 1
                logger.error(f"✗ 解析失败 {file_name}: {e}")
                # 单个文件失败不中断全局流程
                continue
    
    # 输出统计信息
    logger.info(f"\n{'='*50}")
    logger.info(f"批量解析完成")
    logger.info(f"总文件数: {total_files}")
    logger.info(f"成功: {success_files}")
    logger.info(f"失败: {failed_files}")
    logger.info(f"{'='*50}\n")
    
    return results


def get_document_stats(docs: Dict[str, str]) -> dict:
    """
    获取文档统计信息
    
    参数：
        docs: batch_load_docs 返回的字典
    
    返回：
        统计信息字典
    """
    stats = {
        "total_files": len(docs),
        "total_chars": sum(len(content) for content in docs.values()),
        "files_by_type": {},
        "avg_chars_per_file": 0
    }
    
    for file_path, content in docs.items():
        file_type = detect_file_type(file_path)
        if file_type not in stats["files_by_type"]:
            stats["files_by_type"][file_type] = {"count": 0, "chars": 0}
        stats["files_by_type"][file_type]["count"] += 1
        stats["files_by_type"][file_type]["chars"] += len(content)
    
    if stats["total_files"] > 0:
        stats["avg_chars_per_file"] = stats["total_chars"] // stats["total_files"]
    
    return stats


# ==================== 快速使用示例 ====================

if __name__ == "__main__":
    """
    快速使用示例
    """
    import sys
    
    print("=" * 60)
    print("文档解析与加载模块 - Karpathy 知识库架构")
    print("=" * 60)
    
    # 示例 1: 单文件解析
    print("\n【示例 1】单文件解析")
    print("-" * 40)
    
    test_file = "raw/example.txt"
    if os.path.exists(test_file):
        text = extract_full_text(test_file)
        print(f"文件: {test_file}")
        print(f"内容长度: {len(text)} 字符")
        print(f"内容预览: {text[:200]}...")
    else:
        print(f"测试文件不存在: {test_file}")
    
    # 示例 2: 批量文件夹解析
    print("\n【示例 2】批量文件夹解析")
    print("-" * 40)
    
    docs = batch_load_docs("raw/")
    
    if docs:
        stats = get_document_stats(docs)
        print(f"解析文件数: {stats['total_files']}")
        print(f"总字符数: {stats['total_chars']}")
        print(f"平均每文件: {stats['avg_chars_per_file']} 字符")
        print(f"\n文件类型分布:")
        for file_type, info in stats['files_by_type'].items():
            print(f"  {file_type}: {info['count']} 个文件, {info['chars']} 字符")
    else:
        print("未找到可解析的文件")
    
    print("\n" + "=" * 60)
    print("示例完成！")
    print("=" * 60)